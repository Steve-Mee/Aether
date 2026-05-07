#!/usr/bin/env python3
"""
AETHER v1.1 Pitch Deck Generator
Professional 11-slide investor deck
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# Colors
DARK = RGBColor(10, 10, 15)
ACCENT_BLUE = RGBColor(0, 212, 255)
ACCENT_PURPLE = RGBColor(124, 58, 237)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(200, 200, 200)
GREEN = RGBColor(34, 197, 94)

def set_shape_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=12, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, DARK)
    background.line.fill.background()
    
    # Logo
    logo_path = "/home/workdir/artifacts/imagine_images/1HdES.jpg"
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(5.2), Inches(0.8), height=Inches(2.2))
    
    add_text_box(slide, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.8), 
                 "AETHER", 48, True, WHITE, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.5), 
                 "The World's First Self-Evolving AI E-commerce Organism", 18, False, ACCENT_BLUE, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.4), 
                 "Seed Round €2.5 – 3.5M  |  5 mei 2026", 14, False, LIGHT_GRAY, PP_ALIGN.CENTER)
    
    # Slide 2: Problem
    slide = prs.slides.add_slide(slide_layout)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, DARK)
    background.line.fill.background()
    
    add_text_box(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.6), "HET PROBLEEM", 28, True, ACCENT_PURPLE)
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(1.5), 
                 "95%+ van alle nieuwe e-commerce platformen faalt of levert slechts marginale verbeteringen.\nMerchants verliezen 12-20 uur per week aan handmatige admin, email overload en leveranciersbeheer.\nBestaande tools straffen groei af met 2-3% transactiekosten. AI is altijd een 'feature', nooit het DNA.", 
                 14, False, WHITE)
    
    add_text_box(slide, Inches(0.5), Inches(3.5), Inches(12), Inches(0.5), "Resultaat:", 16, True, ACCENT_BLUE)
    add_text_box(slide, Inches(0.5), Inches(4.0), Inches(12), Inches(1.2), 
                 "• Lage winstmarges\n• Hoge churn\n• Constante technische rompslomp\n• Merchants zijn wanhopig op zoek naar échte autonomie", 
                 13, False, LIGHT_GRAY)
    
    # Slide 3: Solution
    slide = prs.slides.add_slide(slide_layout)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, DARK)
    background.line.fill.background()
    
    add_text_box(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.6), "DE OPLOSSING — AETHER", 28, True, ACCENT_BLUE)
    add_text_box(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.8), 
                 "Binnen 60 seconden staat je volledig geoptimaliseerde, wereldwijde webshop live — met AI die 24/7 autonoom:", 
                 14, False, WHITE)
    
    add_text_box(slide, Inches(0.5), Inches(1.8), Inches(4), Inches(1.5), 
                 "AETHER Mail\n70%+ emails autonoom\ngeclassificeerd & beantwoord\n(lokale LLM, zero leakage)", 12, False, LIGHT_GRAY)
    add_text_box(slide, Inches(4.7), Inches(1.8), Inches(4), Inches(1.5), 
                 "Supplier Intelligence\n85%+ accurate prijs/voorraad\nsync via veilige sandbox\n(robots.txt + approval layers)", 12, False, LIGHT_GRAY)
    add_text_box(slide, Inches(8.9), Inches(1.8), Inches(4), Inches(1.5), 
                 "AI-Native Admin\nNatuurlijke taal commando's\n+ 50% minder backend tijd\n(realtime AI inzichten)", 12, False, LIGHT_GRAY)
    
    add_text_box(slide, Inches(0.5), Inches(4.0), Inches(12), Inches(0.5), "UNIEKE SELLING POINTS", 14, True, ACCENT_PURPLE)
    add_text_box(slide, Inches(0.5), Inches(4.5), Inches(12), Inches(0.8), 
                 "✓ 0% Transactiekosten (voor altijd)    ✓ Success-based 12-15% van extra omzet    ✓ Local AI First (air-gapped)    ✓ Volledige Merchant Vrijheid", 
                 12, False, ACCENT_BLUE)
    
    # Slide 4: Roadmap
    slide = prs.slides.add_slide(slide_layout)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, DARK)
    background.line.fill.background()
    
    add_text_box(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.6), "ROADMAP & TRACTION", 28, True, ACCENT_PURPLE)
    
    # Table-like text
    roadmap = """Fase 0 (Q2-Q3 2026)     Foundation + 100 pilots                    —                  100 merchants
Fase 1 (Q4 2026-Q1 2027) AI Brain + Mail + Supplier + Admin        €50M GMV         5.000 merchants
Fase 2 (2027)            Global Scale + Hive Mind                 €2B GMV          50.000 merchants
Fase 3-4 (2028-2029)     Radical Autonomy → Ecosystem Dominance   €100B+ GMV       1M+ merchants"""
    
    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(2.5), roadmap, 11, False, LIGHT_GRAY)
    
    add_text_box(slide, Inches(0.5), Inches(4.2), Inches(12), Inches(0.5), "Huidige status (mei 2026): Volledige v1.1 roadmap + PoC-specificaties + Sprint 1+2+3 code klaar.", 12, False, ACCENT_BLUE)
    
    # Slide 5: Business Model
    slide = prs.slides.add_slide(slide_layout)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, DARK)
    background.line.fill.background()
    
    add_text_box(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.6), "BUSINESS MODEL (Merchant Success First)", 28, True, ACCENT_BLUE)
    
    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(2), 
                 "Pricing\n• Freemium\n• Pro: €99/maand\n• Enterprise: custom\n\nRevenue Streams\n• Success fee 12-15% van extra omzet\n• Data marketplace (Fase 4)\n• 0% platform fees = oneindige schaal", 
                 12, False, LIGHT_GRAY)
    
    add_text_box(slide, Inches(6.5), Inches(1.2), Inches(6), Inches(2), 
                 "Unit Economics (Fase 2 projectie)\n• Gemiddelde uplift: +15-25%\n• Churn target: <8%\n• LTV:CAC > 8:1\n\nWaarom dit werkt:\nMerchants verdienen significant meer →\nwij verdienen alleen als zij méér verdienen.", 
                 12, False, LIGHT_GRAY)
    
    # Slide 6: Team
    slide = prs.slides.add_slide(slide_layout)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, DARK)
    background.line.fill.background()
    
    add_text_box(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.6), "TEAM & MINDSET", 28, True, ACCENT_PURPLE)
    
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(0.8), 
                 "Elon Musk Protocol: First Principles • Radical Simplicity • Boundary Pushing • 100% Honesty • 10x Thinking", 
                 13, True, ACCENT_BLUE)
    
    add_text_box(slide, Inches(0.5), Inches(2.0), Inches(6), Inches(2.5), 
                 "Core Team (Fase 1)\n• CEO / Vision Keeper\n• CTO + Head of AI\n• 18 engineers\n• 7 AI specialists (lokale LLMs)\n• 3 data scientists\n• Hiring: 32 mensen eind 2026", 
                 11, False, LIGHT_GRAY)
    
    add_text_box(slide, Inches(6.5), Inches(2.0), Inches(6), Inches(2.5), 
                 "Waarom dit team wint\n• Diepe expertise in multi-agent systemen\n• Productie-ervaring met lokale LLMs\n• Eerste Principles mindset\n• Merchant Success First als religie", 
                 11, False, LIGHT_GRAY)
    
    # Slide 7: The Ask
    slide = prs.slides.add_slide(slide_layout)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, DARK)
    background.line.fill.background()
    
    add_text_box(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.6), "DE ASK — SEED 2026", 28, True, GREEN)
    
    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.8), "€2.5 – 3.5M Seed Round", 24, True, WHITE)
    
    add_text_box(slide, Inches(0.5), Inches(2.2), Inches(4), Inches(2), 
                 "Gebruik van kapitaal\n60% Team (AI research + engineering)\n25% Lokale GPU infra + security\n(air-gapped agents)\n15% Marketing & eerste 100 pilots", 
                 11, False, LIGHT_GRAY)
    
    add_text_box(slide, Inches(5), Inches(2.2), Inches(7), Inches(2), 
                 "Waarom nu investeren?\n• MedusaJS 2.x is volwassen en extreem uitbreidbaar\n• Lokale LLM modellen (Llama 3.1 70B, Qwen2-VL) zijn eindelijk production-ready\n• Merchants eisen échte autonomie in plaats van meer apps\n• Eerste mover advantage in Autonomous Commerce", 
                 11, False, LIGHT_GRAY)
    
    add_text_box(slide, Inches(0.5), Inches(4.8), Inches(12), Inches(0.5), "Exit potential: Strategische overname (Adobe / Salesforce / Stripe) of IPO 2029+ (€100B+ GMV)", 12, False, ACCENT_BLUE)
    
    # Slide 8: Why AETHER wins
    slide = prs.slides.add_slide(slide_layout)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, DARK)
    background.line.fill.background()
    
    add_text_box(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.6), "WAAROM AETHER WINT", 28, True, ACCENT_PURPLE)
    
    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(4), Inches(2), 
                 "Shopify\nTransactiekosten + app ecosystem\n→ Groei wordt bestraft", 
                 12, False, LIGHT_GRAY)
    
    add_text_box(slide, Inches(4.7), Inches(1.2), Inches(4), Inches(2), 
                 "BigCommerce\nEnterprise bloat\n→ Te complex, te duur", 
                 12, False, LIGHT_GRAY)
    
    add_text_box(slide, Inches(8.9), Inches(1.2), Inches(4), Inches(2), 
                 "WooCommerce\nTechnische rompslomp\n→ Developers nodig", 
                 12, False, LIGHT_GRAY)
    
    add_text_box(slide, Inches(0.5), Inches(3.8), Inches(12), Inches(1), 
                 "AETHER = Merchant Success First + Local AI DNA + 0% fees + Volledige Autonomie\n\nWij maken Shopify irrelevant.", 
                 14, True, ACCENT_BLUE, PP_ALIGN.CENTER)
    
    # Slide 9: Tech Stack Deep Dive
    slide = prs.slides.add_slide(slide_layout)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, DARK)
    background.line.fill.background()
    
    add_text_box(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.6), "TECH STACK (Production-Ready)", 28, True, ACCENT_BLUE)
    
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(6), Inches(3), 
                 "Frontend\n• Next.js 16 + Tailwind + shadcn/ui\n• Medusa Admin extensies (React)\n• Voice + AR interfaces (Fase 2+)\n\nBackend\n• MedusaJS 2.x (commerce core)\n• LangGraph + CrewAI (multi-agent)\n• FastAPI microservices\n\nAI Layer\n• Lokale modellen (Ollama / vLLM)\n• Llama 3.1 70B + Qwen2-VL (vision)\n• Weaviate (vector RAG)", 
                 10, False, LIGHT_GRAY)
    
    add_text_box(slide, Inches(6.5), Inches(1.1), Inches(6), Inches(3), 
                 "Data & Infra\n• PostgreSQL 17 + TimescaleDB\n• Neo4j (graph) + Redis\n• Apache Pulsar (event bus)\n• Kubernetes + Cloudflare Edge\n• Air-gapped Docker containers\n\nSecurity\n• Zero-trust architecture\n• PCI-DSS Level 1 (via partners)\n• GDPR/CCPA native\n• SOC 2 in progress\n• ZK-proofs pilot (Sprint 4)", 
                 10, False, LIGHT_GRAY)
    
    add_text_box(slide, Inches(0.5), Inches(4.5), Inches(12), Inches(0.6), "Eerste Principles: Alles is modulair, lokaal waar privacy telt, en 10x schaalbaar zonder vendor lock-in.", 11, False, ACCENT_PURPLE)
    
    # Slide 10: Competitor Analysis
    slide = prs.slides.add_slide(slide_layout)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, DARK)
    background.line.fill.background()
    
    add_text_box(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.6), "COMPETITOR ANALYSIS", 28, True, ACCENT_PURPLE)
    
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(3.5), 
                 "Shopify: 2-3% transactiekosten + app chaos → Groei wordt bestraft. Merchants haten het.\nBigCommerce: Enterprise bloat, traag, duur. Geen echte AI.\nWooCommerce: Technische rompslomp, developers nodig, geen autonomie.\n\nAETHER differentiatie:\n• 0% fees (nooit transactiekosten)\n• Lokale AI (niet bolt-on, maar DNA)\n• Autonome agents (geen menselijke tussenkomst nodig)\n• Success-based pricing (wij verdienen alleen als merchant méér verdient)\n• Volledige data vrijheid (export zonder straf)", 
                 11, False, LIGHT_GRAY)
    
    add_text_box(slide, Inches(0.5), Inches(5), Inches(12), Inches(0.5), "First Mover Advantage: Niemand bouwt een levend, zelf-evoluerend AI organisme voor commerce.", 12, True, ACCENT_BLUE)
    
    # Slide 11: Financial Projections
    slide = prs.slides.add_slide(slide_layout)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, DARK)
    background.line.fill.background()
    
    add_text_box(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.6), "FINANCIAL PROJECTIONS (Conservatief)", 28, True, ACCENT_BLUE)
    
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(3), 
                 "2026 (Fase 0-1): €0.2M revenue | €1.85M burn | 100 merchants\n2027 (Fase 2): €4.5M revenue | €5.2M burn | 50.000 merchants | Break-even Q4\n2028 (Fase 3): €26M revenue | €12.5M burn | 250.000 merchants | Profitable\n2029 (Fase 4): €125M+ revenue | Self-funding | 1M+ merchants | IPO ready\n\nKey assumptions:\n• Gemiddelde GMV per merchant: €120k (jaar 2)\n• Take rate: 12-15% van extra omzet (niet totale omzet)\n• Churn: <8% (door radicale merchant success)", 
                 10, False, LIGHT_GRAY)
    
    add_text_box(slide, Inches(0.5), Inches(4.5), Inches(12), Inches(0.6), "Path to €1B valuation: 2028 (bij €26M revenue + 30%+ uplift bewezen).", 11, True, GREEN)
    
    # Slide 12: Contact / Closing
    slide = prs.slides.add_slide(slide_layout)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, DARK)
    background.line.fill.background()
    
    # Refined logo
    refined_logo = "/home/workdir/artifacts/imagine_images/6ZYdM.jpg"
    if os.path.exists(refined_logo):
        slide.shapes.add_picture(refined_logo, Inches(5), Inches(0.8), height=Inches(2.5))
    
    add_text_box(slide, Inches(0.5), Inches(3.6), Inches(12), Inches(0.7), "AETHER", 36, True, WHITE, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(4.2), Inches(12), Inches(0.5), "Merchant Success First. Local AI First. Niets is onmogelijk.", 14, False, ACCENT_BLUE, PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(0.5), Inches(5.2), Inches(12), Inches(0.8), 
                 "Steve Meerschaut — Vision Keeper\nGent, België  •  steve@aether.com  •  github.com/Steve-Mee/Aether", 
                 12, False, LIGHT_GRAY, PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(0.5), Inches(6.3), Inches(12), Inches(0.4), 
                 "Seed Round 2026 — €2.5-3.5M — Join us in building the future of commerce.", 
                 11, True, ACCENT_PURPLE, PP_ALIGN.CENTER)
    
    prs.save("/home/workdir/artifacts/AETHER_Pitch_Deck_v1.1.pptx")
    print("Pitch deck created: /home/workdir/artifacts/AETHER_Pitch_Deck_v1.1.pptx")

if __name__ == "__main__":
    create_deck()